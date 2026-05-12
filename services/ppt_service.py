from __future__ import annotations

from datetime import datetime
from pathlib import Path

from loguru import logger
from pptx import Presentation
from pptx.util import Inches

from services.layout_service import LayoutService


EXPORTS_DIR = Path('data/exports')
EXPORTS_DIR.mkdir(parents=True, exist_ok=True)

TEMPLATE_PATH = Path('templates/base_template.pptx')


class PPTService:
    """Production PPTX rendering engine."""

    def __init__(self):
        self.layout_service = LayoutService()
        self.layout = self.layout_service.load_layout()

    def _load_presentation(self) -> Presentation:

        if TEMPLATE_PATH.exists():
            logger.info('Using template PPTX: {}', TEMPLATE_PATH)
            return Presentation(str(TEMPLATE_PATH))

        logger.warning('Template PPTX not found. Using fallback presentation.')
        return Presentation()

    def generate_presentation(
        self,
        period_label: str,
        total_calls: int,
        topics_chart_path: str | None = None,
        top5_chart_path: str | None = None,
        dynamics_chart_path: str | None = None
    ) -> str:

        presentation = self._load_presentation()

        presentation.slide_width = Inches(
            self.layout['slide']['width']
        )

        presentation.slide_height = Inches(
            self.layout['slide']['height']
        )

        slide_layout = presentation.slide_layouts[6]
        slide = presentation.slides.add_slide(slide_layout)

        title_layout = self.layout['title']

        title_box = slide.shapes.add_textbox(
            Inches(title_layout['x']),
            Inches(title_layout['y']),
            Inches(title_layout['w']),
            Inches(title_layout['h'])
        )

        title_frame = title_box.text_frame

        title_frame.text = (
            'ДИНАМИКА ОБРАЩЕНИЙ ПОТРЕБИТЕЛЕЙ '
            f'ПО ЭЗС ЗА {period_label}'
        )

        metrics_box = slide.shapes.add_textbox(
            Inches(0.3),
            Inches(0.75),
            Inches(4),
            Inches(0.4)
        )

        metrics_frame = metrics_box.text_frame
        metrics_frame.text = f'Обращений: {total_calls}'

        if topics_chart_path:

            chart_layout = self.layout['topics_chart']

            slide.shapes.add_picture(
                topics_chart_path,
                Inches(chart_layout['x']),
                Inches(chart_layout['y']),
                Inches(chart_layout['w']),
                Inches(chart_layout['h'])
            )

        if top5_chart_path:

            chart_layout = self.layout['top5_chart']

            slide.shapes.add_picture(
                top5_chart_path,
                Inches(chart_layout['x']),
                Inches(chart_layout['y']),
                Inches(chart_layout['w']),
                Inches(chart_layout['h'])
            )

        if dynamics_chart_path:

            chart_layout = self.layout['dynamics_chart']

            slide.shapes.add_picture(
                dynamics_chart_path,
                Inches(chart_layout['x']),
                Inches(chart_layout['y']),
                Inches(chart_layout['w']),
                Inches(chart_layout['h'])
            )

        generated_at = datetime.now().strftime('%Y%m%d_%H%M%S')

        output_path = (
            EXPORTS_DIR /
            f'ev_charging_report_{generated_at}.pptx'
        )

        presentation.save(output_path)

        logger.info(
            'Presentation exported successfully: {}',
            output_path
        )

        return str(output_path)
